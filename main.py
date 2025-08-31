from typing import TypedDict, List, Optional, Annotated, Literal
from langgraph.graph import StateGraph, END, add_messages
import json
import os
import requests
import time
import subprocess
import uuid
import python_dotenv
# --- Files
import pymupdf  # PyMuPDF
import docx
import pandas as pd
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter



# --- LLM & messages
from langchain_openai import ChatOpenAI
from langchain_core.messages import SystemMessage, HumanMessage, AIMessage
from langchain_core.prompts import PromptTemplate

# --- Vector DB
import chromadb
from chromadb.config import Settings
from sentence_transformers import SentenceTransformer

# --- Tools
from langchain.tools import tool
from langchain.chains import LLMMathChain
from langchain_tavily import TavilySearch
from langchain_experimental.utilities import PythonREPL
from langchain.agents import AgentExecutor, create_react_agent

# -------------- KEYS / ENV --------------


OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY")
ASSEMBLYAI_API_KEY = os.environ.get("ASSEMBLYAI_API_KEY")
TAVILY_API_KEY = os.environ.get("TAVILY_API_KEY")

splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)

# -------------- LLM INIT --------------
llm = ChatOpenAI(
    model_name="gpt-4o-mini",   
    temperature=0.0,
    openai_api_key=OPENAI_API_KEY,
)

# -------------- STATE --------------
class State(TypedDict):
    content_id: str
    file_type: Optional[str]
    raw_content: Optional[str]
    processed_content: Optional[str]
    summary: Optional[str]
    study_aids: Optional[dict]
    study_aids_preference: Optional[Literal["quizzes", "flashcards", "both", "none"]]
    question: Optional[str]
    response: Optional[str]
    feedback: Optional[int]
    needs_reprocessing: Optional[bool]
    messages: Annotated[List, add_messages]

# -------------- TRANSCRIPTION --------------
class AssemblyAPI:
    def __init__(self, api_key: Optional[str]):
        self.api_key = api_key
        self.headers = {
            "authorization": self.api_key or "",
            "content-type": "application/json"
        }
        self.base_url = "https://api.assemblyai.com/v2"

    def _upload_file(self, file_path):
        if not self.api_key:
            raise RuntimeError("ASSEMBLYAI_API_KEY not set.")
        with open(file_path, "rb") as f:
            
            response = requests.post(
                f"{self.base_url}/upload",
                headers={"authorization": self.api_key},
                files={"file": f}
            )
        response.raise_for_status()
        return response.json()["upload_url"]

    def _submit_transcription(self, audio_url):
        if not self.api_key:
            raise RuntimeError("ASSEMBLYAI_API_KEY not set.")
        response = requests.post(
            f"{self.base_url}/transcript",
            json={"audio_url": audio_url},
            headers=self.headers
        )
        response.raise_for_status()
        return response.json()["id"]

    def _poll_transcription(self, transcript_id):
        polling_url = f"{self.base_url}/transcript/{transcript_id}"
        while True:
            response = requests.get(polling_url, headers=self.headers)
            response.raise_for_status()
            data = response.json()
            if data["status"] == "completed":
                return data["text"]
            if data["status"] == "error":
                raise Exception(f"Transcription failed: {data['error']}")
            time.sleep(3)

    def transcribe_youtube(self, youtube_url):
        file_name = f"{uuid.uuid4()}.mp3"
        try:
            subprocess.run(
                ["yt-dlp", "--extract-audio", "--audio-format", "mp3", "-o", file_name, youtube_url],
                check=True
            )
            upload_url = self._upload_file(file_name)
            transcript_id = self._submit_transcription(upload_url)
            return self._poll_transcription(transcript_id)
        finally:
            if os.path.exists(file_name):
                os.remove(file_name)

    def transcribe_audio_file(self, audio_file_path):
        upload_url = self._upload_file(audio_file_path)
        transcript_id = self._submit_transcription(upload_url)
        return self._poll_transcription(transcript_id)


assembly_api = AssemblyAPI(api_key=ASSEMBLYAI_API_KEY)

# -------------- CHROMA WRAPPER --------------
class ChromaVectorDB:
    def __init__(self, collection_name="study_aid_collection"):
        self.embedding_model = SentenceTransformer("all-MiniLM-L6-v2")
        self.chroma_client = chromadb.Client(Settings(anonymized_telemetry=False))
        self.collection = self.chroma_client.get_or_create_collection(name=collection_name)

    def embed(self, texts: List[str]):
        return self.embedding_model.encode(texts).tolist()

    def store(self, content_id: str, content: str):
        embedding = self.embed([content])[0]
        self.collection.add(documents=[content], embeddings=[embedding], ids=[content_id])
        return True

    def search(self, query: str, n_results: int = 3):
        embedding = self.embed([query])[0]
        results = self.collection.query(query_embeddings=[embedding], n_results=n_results)
        hits = []
        for doc, doc_id in zip(results.get("documents", [[]])[0], results.get("ids", [[]])[0]):
            hits.append({"content_id": doc_id, "content": (doc[:400] + "...") if doc else ""})
        return hits

vectordb = ChromaVectorDB()

# -------------- LLM HELPERS --------------
def llm_summarize(text: str, max_tokens: int = 512) -> str:
    prompt = [
        SystemMessage(content="You are an assistant that writes concise summaries."),
        HumanMessage(content=f"Summarize the following text in 3-6 bullet points:\n\n{text[:15000]}")
    ]
    resp = llm.invoke(prompt)
    return resp.content if hasattr(resp, "content") else str(resp)

def llm_generate_study_aids(text: str, preference: str = "both") -> dict:
    prompt = [
        SystemMessage(content="You are an assistant that creates study aids (quizzes, flashcards) from a text."),
        HumanMessage(content=f"Create study aids (preference={preference}) for the following content:\n\n{text[:8000]}")
    ]
    try:
        resp = llm.invoke(prompt)
        out = resp.content if hasattr(resp, "content") else str(resp)
    except Exception as e:
        out = f"Failed to generate study aids: {e}"
    return {"raw": out} 

# -------------- GRAPH NODES --------------
def router_node(state: State):
    """Route based on input type"""
    if state.get("question"):
        return {"question": state["question"]}
    if state.get("file_type"):
        return {"file_type": state["file_type"]}
    raise ValueError("No question or file type provided")

def transcribe_node(state: State):
    """Transcribe YouTube or audio file using AssemblyAI"""
    file_type = state.get("file_type")
    raw_content = state.get("raw_content")
    if not raw_content:
        return {"processed_content": ""}
    if file_type == "video":
        transcription = assembly_api.transcribe_youtube(raw_content)
    elif file_type == "audio":
        transcription = assembly_api.transcribe_audio_file(raw_content)
    else:
        transcription = ""
    return {"processed_content": transcription}

def parse_document_node(state):
    """Parse different document types into text"""
    file_path = state.get("raw_content", "")
    processed_text = ""

    if not file_path or not os.path.exists(file_path):
        return {"processed_content": ""}

    ext = os.path.splitext(file_path)[1].lower()

    try:
        # PDF
        if ext == ".pdf":
            doc = pymupdf.open(file_path)
            for page in doc:
                processed_text += page.get_text()
            doc.close()

        # DOCX
        elif ext == ".docx":
            doc = docx.Document(file_path)
            processed_text = "\n".join([para.text for para in doc.paragraphs])

        # CSV
        elif ext == ".csv":
            df = pd.read_csv(file_path)
            processed_text = df.to_string()

        # PPTX
        elif ext == ".pptx":
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        processed_text += shape.text + "\n"

        # TXT (fallback)
        elif ext == ".txt":
            with open(file_path, "r", encoding="utf-8") as f:
                processed_text = f.read()

        else:
            processed_text = f"Unsupported file type: {ext}"

    except Exception as e:
        processed_text = f"Error parsing file: {str(e)}"

    return {"processed_content": processed_text[:1000]}  # limit length


def aggregate_content_node(state: State):
    """Placeholder: return state untouched (only partial updates recommended in practice)."""
    return {}

def generate_summary_node(state: State):
    """Generate summary from processed/raw content using LLM"""
    text = state.get("processed_content") or state.get("raw_content") or ""
    if not text:
        return {"summary": ""}
    return {"summary": llm_summarize(text)}

def study_aids_preference_node(state: State):
    preference = state.get("study_aids_preference", "both")
    return {"study_aids_preference": preference}

def generate_study_aids_node(state: State):
    content = state.get("processed_content", "")
    preference = state.get("study_aids_preference", "both")
    if not content:
        return {"study_aids": None}
    return {"study_aids": llm_generate_study_aids(content, preference=preference)}

def index_vectordb_node(state: State):
    content = state.get("processed_content")
    if not content:
        return {}

    chunks = splitter.split_text(content)
    content_id = state.get("content_id") or str(uuid.uuid4())

    for i, chunk in enumerate(chunks):
        vectordb.store(f"{content_id}_{i}", chunk)

    return {"vectordb_status": f"Indexed {len(chunks)} chunks under {content_id}"}

# ---------------- REAL TOOLS ----------------
# 1) Calculator (LLM MathChain)
math_chain = LLMMathChain.from_llm(llm)

@tool
def calculator_tool(question: str) -> str:
    """Use when you need to do math or evaluate expressions. Input: the math question."""
    try:
        return math_chain.invoke(question)
    except Exception as e:
        return f"Math error: {e}"

# 2) Web search (Tavily)
tavily_tool = TavilySearch(max_results=3)

@tool
def web_search_tool(query: str) -> str:
    """Search the web for current/recent info. Input: search query."""
    try:
        results = tavily_tool.invoke({"query": query})  
        # If tavily_tool returns a dict instead of list, handle accordingly
        if isinstance(results, dict) and "results" in results:
            results = results["results"]

        lines = []
        for r in results:
            title = r.get("title") or "Result"
            url = r.get("url") or ""
            snippet = (r.get("content") or "")[:220]
            lines.append(f"- {title}\n  {url}\n  {snippet}")

        return "\n".join(lines) if lines else "No results."
    except Exception as e:
        return f"Search error: {e}"


# 3) Code execution (Python REPL)
python_repl = PythonREPL()

@tool
def code_execution_tool(code: str) -> str:
    """Run short Python code snippets. Input: Python code."""
    try:
        return str(python_repl.run(code))   
    except Exception as e:
        return f"Code error: {e}"

# 4) RAG (Chroma search)

@tool
def basic_rag(query: str) -> str:
    """Query the internal knowledge base. Input: your question."""
    hits = vectordb.search(query, n_results=3)
    if not hits:
        return "No relevant information found in the knowledge base."

    out = []
    for h in hits:
        snippet = h['content'][:300] + "..." if len(h['content']) > 300 else h['content']
        source = h.get("source", "unknown")
        score = h.get("score", None)
        if score:
            out.append(f"- [{h['content_id']}] {snippet} (source: {source}, score: {score:.2f})")
        else:
            out.append(f"- [{h['content_id']}] {snippet} (source: {source})")

    return "\n".join(out)

tools = [calculator_tool, web_search_tool, code_execution_tool, basic_rag]
llm_with_tools = llm.bind_tools(tools)



template = """You are a reasoning agent that can use tools.
{tools}
Use the following format:

Question: the input question
Thought: you should always think about what to do
Action: the action to take (one of {tool_names})
Action Input: the input to the action
Observation: the result of the action
... (this Thought/Action/Action Input/Observation can repeat N times)
Final Answer: the final answer to the original input question

Begin!

Question: {input}
{agent_scratchpad}
"""

prompt = PromptTemplate(
    input_variables=["input", "agent_scratchpad"],
    template=template,
)

# ---- Create ReAct Agent ----
agent = create_react_agent(llm, tools, prompt)


agent_executor = AgentExecutor.from_agent_and_tools(
    agent=agent,
    tools=tools,
    handle_parsing_errors=True,
    verbose=True
)

def call_model(state: State):
    """Run the agent executor with question/messages as input."""
    msgs = state.get("messages", [])
    question = state.get("question")

    if question:
        msgs = msgs + [HumanMessage(content=question)]

    
    result = agent_executor.invoke({"input": question or "", "chat_history": msgs})

    return {
        "response": result.get("output", str(result)),
        "messages": msgs,
    }


# -------------- FORMAT & FEEDBACK --------------
def format_response_node(state: State):
    parts = []
    if summary := state.get("summary"):
        parts.append(f"Summary:\n{summary}")
    if study_aids := state.get("study_aids"):
        try:
            study_text = json.dumps(study_aids, indent=2) if isinstance(study_aids, (dict, list)) else str(study_aids)
        except Exception:
            study_text = str(study_aids)
        parts.append(f"Study Aids:\n{study_text}")
    if response := state.get("response"):
        parts.append(f"Response:\n{response}")
    return {"response": "\n\n".join(parts) if parts else "No response generated"}

def process_feedback_node(state: State):
    feedback = state.get("feedback")
    needs_reprocessing = (feedback is not None) and (feedback < 3)
    return {"needs_reprocessing": needs_reprocessing, "feedback": feedback}



# -------------- GRAPH --------------
builder = StateGraph(State)

# nodes
builder.add_node("router", router_node)
builder.add_node("transcribe", transcribe_node)
builder.add_node("parse_document", parse_document_node)
builder.add_node("aggregate_content", aggregate_content_node)
builder.add_node("generate_summary", generate_summary_node)
builder.add_node("study_aids_preference", study_aids_preference_node)
builder.add_node("generate_study_aids", generate_study_aids_node)
builder.add_node("index_vectordb", index_vectordb_node)
builder.add_node("call_model", call_model)   # renamed correctly
builder.add_node("format_response", format_response_node)
builder.add_node("process_feedback", process_feedback_node)

# entry
builder.set_entry_point("router")

# branching logic in router
def router_condition(state: State):
    if state.get("question"):   # Q&A branch
        return "call_model"
    if state.get("file_type") in {"video", "audio"}:
        return "transcribe"
    if state.get("file_type"):
        return "parse_document"
    raise ValueError("No question or file type provided")

builder.add_conditional_edges(
    "router",
    router_condition,
    ["call_model", "transcribe", "parse_document"]
)

# content processing branch
builder.add_edge("transcribe", "aggregate_content")
builder.add_edge("parse_document", "aggregate_content")

# from aggregate_content → both summary & vectordb
builder.add_edge("aggregate_content", "generate_summary")
builder.add_edge("aggregate_content", "index_vectordb")

# after summaries → study aids flow
builder.add_edge("generate_summary", "study_aids_preference")
#builder.add_conditional_edges(
#    "study_aids_preference",
#    lambda s: "generate_study_aids" if s.get("study_aids_preference") != "none" else "index_vectordb",
#    ["generate_study_aids", "index_vectordb"]
#)
#builder.add_edge("generate_study_aids", "index_vectordb")

# vectordb → call_model (RAG Q&A)
#builder.add_edge("index_vectordb", "call_model")
builder.add_edge("call_model", "index_vectordb")

# converge to formatting
builder.add_edge("call_model", "format_response")

builder.add_edge("study_aids_preference", "generate_study_aids")
# feedback loop
builder.add_edge("format_response", "process_feedback")
builder.add_conditional_edges(
    "process_feedback",
    lambda s: "generate_summary" if s.get("needs_reprocessing") else END,
    ["generate_summary", END]
)

graph = builder.compile()
